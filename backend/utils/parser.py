"""
Mimir — Document Parser
Extracts text from PDFs and images, then indexes chunks into ChromaDB.

Requirements (installed separately):
  pip install pymupdf pytesseract pillow
  # Also install Tesseract binary: https://github.com/UB-Mannheim/tesseract/wiki

Graceful degradation: if either library is absent, that path returns empty
text and the file is still marked processed (no crash).
"""

import re
import logging

from sqlalchemy import select

from memory.database import AsyncSessionLocal, File as FileModel, UserMemory, ExamQuestion as ExamQuestionModel
from memory.vector import add_document_memory
from utils.exam_parser import is_exam_paper, parse_exam_questions

logger = logging.getLogger("mimir.parser")

# ── Optional imports — degrade gracefully if not installed ───
try:
    import fitz  # PyMuPDF
    _HAS_PYMUPDF = True
except ImportError:
    _HAS_PYMUPDF = False
    logger.warning("PyMuPDF not installed — PDF text extraction unavailable.")

try:
    from pptx import Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False
    logger.warning("python-pptx not installed — PowerPoint extraction unavailable.")

try:
    import pytesseract
    from PIL import Image
    # Point pytesseract at the standard Windows install path if not in PATH
    import shutil, os
    if not shutil.which("tesseract"):
        _DEFAULT_TESS = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(_DEFAULT_TESS):
            pytesseract.pytesseract.tesseract_cmd = _DEFAULT_TESS
    _HAS_OCR = True
except ImportError:
    _HAS_OCR = False
    logger.warning("pytesseract / Pillow not installed — image OCR unavailable.")


# ── Text chunking ────────────────────────────────────────────

_SENT_END_RE = re.compile(r'(?<=[.!?])\s+')


def _last_sentence(text: str) -> str:
    """Return the final sentence of *text* (used for chunk overlap)."""
    parts = _SENT_END_RE.split(text.strip())
    return parts[-1].strip() if parts else ""


def _semantic_chunk(text: str, max_size: int = 800) -> list[str]:
    """Paragraph-aware semantic chunker — replaces the old fixed-size splitter.

    Strategy:
    1. Normalise line-endings and collapse excessive blank lines.
    2. Split on blank lines (paragraph boundaries).
    3. Merge adjacent short paragraphs up to *max_size*.
    4. Split long paragraphs on sentence boundaries.
    5. Overlap: include the last sentence of each emitted chunk at the start
       of the next chunk so ChromaDB retrieval spans boundaries cleanly.

    Args:
        text:     Raw extracted text (may contain multiple blank lines).
        max_size: Target maximum chunk size in characters.

    Returns:
        List of non-empty chunk strings.
    """
    # Normalise
    text = re.sub(r'\r\n?', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)

    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    if not paragraphs:
        return []

    chunks: list[str] = []
    buf  = ""
    tail = ""   # last sentence of the most recently emitted chunk

    def flush(b: str) -> None:
        nonlocal tail
        b = b.strip()
        if len(b) > 15:
            chunks.append(b)
            tail = _last_sentence(b)

    for para in paragraphs:
        if len(para) > max_size:
            # Flush current buffer first, then split the long paragraph
            if buf:
                flush(buf)
                buf = ""

            sents = _SENT_END_RE.split(para)
            seg = tail  # begin with overlap from previous chunk
            for s in sents:
                s = s.strip()
                if not s:
                    continue
                candidate = (seg + " " + s).strip() if seg else s
                if len(candidate) > max_size and seg:
                    flush(seg)
                    seg = (tail + " " + s).strip() if tail else s
                else:
                    seg = candidate
            if seg:
                buf = seg   # keep remainder; may merge with next paragraph
        else:
            candidate = (buf + "\n\n" + para).strip() if buf else para
            if len(candidate) > max_size and buf:
                flush(buf)
                buf = (tail + "\n\n" + para).strip() if tail else para
            else:
                buf = candidate

    if buf:
        flush(buf)

    return chunks


# ── Extractors ───────────────────────────────────────────────
def _extract_pdf(filepath: str) -> str:
    """Extract text from a PDF preserving document structure using PyMuPDF.

    Uses block-level extraction with font-size heuristics to detect headings.
    Headings are prefixed with a blank line so the semantic chunker treats them
    as paragraph boundaries, keeping related content together.

    Falls back to plain ``get_text()`` if structured extraction fails.
    """
    if not _HAS_PYMUPDF:
        return ""
    try:
        doc = fitz.open(filepath)
        page_texts: list[str] = []

        for page in doc:
            try:
                raw = page.get_text("dict", sort=True)
                blocks = raw.get("blocks", [])
                parts: list[str] = []

                for block in blocks:
                    if block.get("type") != 0:   # skip image blocks
                        continue

                    # Collect all spans in the block to measure font sizes
                    spans: list[dict] = [
                        span
                        for line in block.get("lines", [])
                        for span in line.get("spans", [])
                    ]
                    if not spans:
                        continue

                    block_text = " ".join(s.get("text", "") for s in spans).strip()
                    if not block_text:
                        continue

                    # Heading heuristic: max font size ≥ 13 AND short line
                    max_size = max(s.get("size", 0) for s in spans)
                    is_heading = max_size >= 13 and len(block_text) <= 150

                    if is_heading:
                        # Blank line before heading creates a semantic boundary
                        parts.append(f"\n{block_text}")
                    else:
                        parts.append(block_text)

                page_texts.append("\n\n".join(parts))

            except Exception:
                # Per-page fallback to plain text
                page_texts.append(page.get_text())

        doc.close()
        return "\n\n".join(page_texts)

    except Exception as e:
        logger.error("PDF extraction failed for %s: %s", filepath, e)
        return ""


def _extract_pptx(filepath: str) -> str:
    """Extract text from a PowerPoint deck, one section per slide.

    Lecture slides are how most course material actually arrives, so they are
    worth reading properly rather than asking students to convert to PDF first.

    Slide titles are emitted on their own line with a blank line before them, so
    the semantic chunker treats each slide as a boundary and never merges the
    end of one slide into the start of the next — the same trick the PDF path
    uses with font-size headings.

    Speaker notes are included: in a lecture deck they frequently carry the
    explanation that the slide itself only gestures at, which is exactly the
    material a student later wants to search.

    Only the modern ``.pptx`` format is readable. The legacy binary ``.ppt`` is
    a different container that python-pptx cannot open; the router rejects it
    with a message telling the user to re-save.
    """
    if not _HAS_PPTX:
        return ""
    try:
        prs = Presentation(filepath)
        slides: list[str] = []

        for number, slide in enumerate(prs.slides, start=1):
            title = ""
            body: list[str] = []

            for shape in slide.shapes:
                # Tables carry real content in technical decks.
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            body.append(" | ".join(cells))
                    continue

                if not getattr(shape, "has_text_frame", False):
                    continue

                text = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                ).strip()
                if not text:
                    continue

                # The placeholder marked as the title, when the layout has one.
                is_title = False
                try:
                    ph = getattr(shape, "placeholder_format", None)
                    is_title = ph is not None and ph.idx == 0
                except Exception:
                    is_title = False

                if is_title and not title:
                    title = text.replace("\n", " ")
                else:
                    body.append(text)

            notes = ""
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes = ""

            if not (title or body or notes):
                continue                      # purely decorative slide

            parts = [f"\nSlide {number}: {title}" if title else f"\nSlide {number}"]
            parts.extend(body)
            if notes:
                parts.append(f"Speaker notes: {notes}")
            slides.append("\n".join(parts))

        return "\n\n".join(slides)

    except Exception as e:
        logger.error("PPTX extraction failed for %s: %s", filepath, e)
        return ""


def _extract_image_ocr(filepath: str) -> str:
    """Run Tesseract OCR on an image file. Returns empty string if unavailable."""
    if not _HAS_OCR:
        return ""
    try:
        img = Image.open(filepath)
        return pytesseract.image_to_string(img)
    except Exception as e:
        logger.error("OCR failed for %s: %s", filepath, e)
        return ""


async def _extract_image(filepath: str) -> str:
    """Extract content from an image using the Ollama vision model.

    Strategy:
    1. Try the configured vision model (``settings.vision_model``).
       Produces a rich, structured description of diagrams, text, and
       relationships — far better than raw OCR for study indexing.
    2. Fall back to Tesseract OCR if the vision model is not loaded or
       the call fails.
    3. Return empty string if both fail — the file is still marked
       processed so the upload does not block the user.
    """
    import base64
    try:
        import ollama as _ollama
        from config import settings as _s

        with open(filepath, "rb") as fh:
            img_b64 = base64.b64encode(fh.read()).decode("utf-8")

        client = _ollama.AsyncClient(host=_s.ollama_base_url)
        resp = await client.chat(
            model=_s.vision_model,
            messages=[{
                "role":    "user",
                "content": (
                    "You are indexing this image for a student's study notes. "
                    "Extract and describe ALL content: every text label, formula, "
                    "arrow, relationship, diagram structure, and key concept. "
                    "Be thorough so this image can be found by relevant search queries."
                ),
                "images":  [img_b64],
            }],
            options={"temperature": 0.2},
            stream=False,
        )
        text = resp["message"]["content"].strip()
        if text:
            return text
        logger.warning("Vision model returned empty content for %s — falling back to OCR", filepath)
    except Exception as e:
        logger.warning("Vision model unavailable for %s (%s) — falling back to OCR", filepath, e)

    return _extract_image_ocr(filepath)


async def _store_document_memory(
    user_id: int,
    subject_id: int | None,
    filename: str,
    text_preview: str,
) -> None:
    """Call LLM to extract a structured document summary into user_memories.

    Produces a short summary + key topics so the agent always knows what
    documents the student has uploaded, even without a similarity hit.
    """
    try:
        import ollama as _ollama
        from config import settings as _s

        client = _ollama.AsyncClient(host=_s.ollama_base_url)
        resp = await client.chat(
            model=_s.ollama_model,
            messages=[{
                "role": "user",
                "content": (
                    f"Analyse this document excerpt and respond in EXACTLY this format:\n"
                    f"SUMMARY: <one sentence describing the document>\n"
                    f"TOPICS: <comma-separated key concepts>\n"
                    f"KEY_FACTS: <up to 5 bullet points of important formulas, definitions, or facts>\n\n"
                    f"Document ({filename}):\n{text_preview}"
                ),
            }],
            options={"temperature": 0.2, "num_predict": 300},
            think=False,
        )
        content = resp["message"]["content"].strip()
        if not content:
            return

        async with AsyncSessionLocal() as db:
            mem = UserMemory(
                user_id=user_id,
                subject_id=subject_id,
                memory_type="document",
                key=f"doc:{filename}",
                value=content,
                source=filename,
            )
            db.add(mem)
            await db.commit()
        logger.info("Stored document memory for file %r (user %d)", filename, user_id)
    except Exception as exc:
        logger.warning("Document memory extraction failed for %r: %s", filename, exc)


# ── Main async task ──────────────────────────────────────────
async def parse_and_index_file(
    file_id:      int,
    filepath:     str,
    user_id:      int,
    subject_id:   int | None,
    content_type: str,
    filename:     str = "",
) -> None:
    """
    Background task: extract text → chunk → store in ChromaDB → mark processed.

    Called by FastAPI BackgroundTasks immediately after a successful upload.
    Runs in the same process, after the HTTP response has been sent.
    """
    logger.info("Parsing file %s (id=%d, type=%s)", filepath, file_id, content_type)

    # ── Extract ──────────────────────────────────────────────
    if "pdf" in content_type:
        text = _extract_pdf(filepath)
    elif "presentation" in content_type or filepath.lower().endswith(".pptx"):
        text = _extract_pptx(filepath)
    else:
        text = await _extract_image(filepath)

    # ── Index chunks ─────────────────────────────────────────
    chunk_count = 0
    if text.strip():
        chunks = _semantic_chunk(text)
        chunk_count = len(chunks)
        logger.info("Indexing %d semantic chunks for file %d", chunk_count, file_id)
        for idx, chunk in enumerate(chunks):
            try:
                add_document_memory(
                    user_id=user_id,
                    content=chunk,
                    file_id=file_id,
                    chunk_idx=idx,
                    subject_id=subject_id,
                    filename=filename,
                )
            except Exception as e:
                logger.error("Failed to index chunk %d: %s", idx, e)
    else:
        logger.warning("No text extracted from file %d — check library install.", file_id)

    # ── Exam-paper detection & question extraction ───────────
    exam_question_count = 0
    # Slide decks are checked too — tutorial and revision decks often carry
    # past-paper questions with their mark allocations.
    _is_document = "pdf" in content_type or "presentation" in content_type
    if _is_document and text.strip() and is_exam_paper(text):
        parsed_qs = parse_exam_questions(text, filename)
        if parsed_qs:
            async with AsyncSessionLocal() as db:
                for pq in parsed_qs:
                    db.add(ExamQuestionModel(
                        file_id=file_id,
                        user_id=user_id,
                        subject_id=subject_id,
                        question_number=pq.question_number,
                        question_text=pq.question_text,
                        marks=pq.marks,
                        page_number=pq.page_number,
                    ))
                await db.commit()
            exam_question_count = len(parsed_qs)
            logger.info(
                "Stored %d exam questions from file %d (%r).",
                exam_question_count, file_id, filename,
            )

    # ── Mark processed ───────────────────────────────────────
    async with AsyncSessionLocal() as db:
        result = await db.execute(
            select(FileModel).where(FileModel.id == file_id)
        )
        file_row = result.scalar_one_or_none()
        if file_row:
            file_row.processed = True
            if exam_question_count > 0:
                file_row.has_exam_questions = True
            await db.commit()
            logger.info("File %d marked as processed.", file_id)

    # ── Store document memory (always-on context for agent) ──
    if text.strip():
        await _store_document_memory(user_id, subject_id, filename, text[:3000])

    # ── Notify user via WebSocket ─────────────────────────────
    try:
        from ws_manager import manager
        await manager.send_to_user(user_id, {
            "type":               "file_indexed",
            "file_id":            file_id,
            "filename":           filename,
            "chunks":             chunk_count,
            "exam_question_count": exam_question_count,
        })
        logger.info("Sent file_indexed WS notification to user %d.", user_id)
    except Exception as e:
        logger.warning("Could not send file_indexed WS notification: %s", e)
